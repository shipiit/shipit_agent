"""File-extraction backends — turn an office/web document into clean text.

``read_file`` reads source, config, and logs as UTF-8 and images/PDFs as vision.
Everything else — a ``.docx``, ``.xlsx``, ``.pptx``, an ``.html`` page — is
binary or markup that decodes to U+FFFD soup. This registry fills that gap: one
small extractor per format, each turning a file into **clean Markdown the model
can actually read** (headings, tables, slide breaks preserved).

Same shape as the media registries: each backend declares the suffixes it owns
and a cheap ``is_available()`` (its parsing library present, no I/O), so a format
lights up only when its dependency is installed and otherwise degrades to the
plain-text path. New formats register with :func:`register_extractor` — the same
shape a plugin uses.

Built-in: **CSV/TSV** (stdlib, always on), **DOCX** (python-docx), **XLSX**
(openpyxl), **PPTX** (python-pptx), **HTML** (markdownify, or a stdlib fallback),
**PDF** (pypdf). Install the set with ``pip install 'shipit-agent[files]'``.
"""

from __future__ import annotations

import csv
import importlib.util
from pathlib import Path
from typing import Any, Protocol, runtime_checkable


@runtime_checkable
class Extractor(Protocol):
    name: str
    suffixes: tuple[str, ...]

    def is_available(self) -> bool:
        """Cheap check — the parsing library is importable. No file I/O."""

    def extract(self, path: Path) -> str:
        """Return Markdown/plain text for ``path``. Raises on a real failure."""


_REGISTRY: dict[str, Extractor] = {}


def register_extractor(extractor: Extractor) -> Extractor:
    """Add (or override) an extractor, keyed by ``extractor.name``."""
    _REGISTRY[extractor.name] = extractor
    return extractor


def _has(module: str) -> bool:
    return importlib.util.find_spec(module) is not None


def extractor_for(suffix: str) -> Extractor | None:
    """The first *available* extractor that owns ``suffix`` (e.g. ``.docx``)."""
    suffix = suffix.lower()
    for extractor in _REGISTRY.values():
        if suffix in extractor.suffixes and _safe_available(extractor):
            return extractor
    return None


def extract_file(path: Path) -> tuple[str, str] | None:
    """Extract ``path`` to ``(markdown, extractor_name)``, or ``None``.

    ``None`` means no available extractor owns this suffix — the caller falls
    back to reading the bytes as text. A parsing failure is *not* swallowed
    here; it raises, so the caller can report *why* a supported file failed.
    """
    extractor = extractor_for(path.suffix)
    if extractor is None:
        return None
    return extractor.extract(path), extractor.name


def supported_suffixes() -> set[str]:
    """Every suffix any registered extractor claims (ignores availability)."""
    return {suffix for e in _REGISTRY.values() for suffix in e.suffixes}


def available_suffixes() -> set[str]:
    """Suffixes whose extractor can run right now (dependency installed)."""
    return {
        suffix
        for e in _REGISTRY.values()
        if _safe_available(e)
        for suffix in e.suffixes
    }


def _safe_available(extractor: Extractor) -> bool:
    try:
        return bool(extractor.is_available())
    except Exception:  # noqa: BLE001 — a broken probe means "unavailable"
        return False


def _rows_to_markdown(rows: list[list[Any]], *, max_rows: int = 500) -> str:
    """Render tabular rows as a GitHub-flavoured Markdown table."""
    rows = [[("" if cell is None else str(cell)) for cell in row] for row in rows if row]
    if not rows:
        return "*(empty)*"
    truncated = len(rows) > max_rows
    rows = rows[:max_rows]
    width = max(len(row) for row in rows)
    rows = [row + [""] * (width - len(row)) for row in rows]
    header, *body = rows
    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join(["---"] * width) + " |",
        *("| " + " | ".join(row) + " |" for row in body),
    ]
    if truncated:
        lines.append(f"\n*…{max_rows}-row cap reached; more rows not shown.*")
    return "\n".join(lines)


# ── built-in extractors ──────────────────────────────────────────────────────


class CsvExtractor:
    """CSV / TSV → a Markdown table. Stdlib only — always available."""

    name = "csv"
    suffixes = (".csv", ".tsv")

    def is_available(self) -> bool:
        return True

    def extract(self, path: Path) -> str:
        delimiter = "\t" if path.suffix.lower() == ".tsv" else ","
        with path.open("r", encoding="utf-8", errors="replace", newline="") as handle:
            rows = list(csv.reader(handle, delimiter=delimiter))
        return _rows_to_markdown(rows)


class DocxExtractor:
    """Word ``.docx`` → Markdown (paragraphs, headings, and tables)."""

    name = "docx"
    suffixes = (".docx",)

    def is_available(self) -> bool:
        return _has("docx")

    def extract(self, path: Path) -> str:
        import docx  # python-docx

        document = docx.Document(str(path))
        parts: list[str] = []
        for paragraph in document.paragraphs:
            text = paragraph.text.strip()
            if not text:
                continue
            style = (paragraph.style.name or "").lower() if paragraph.style else ""
            if style.startswith("heading"):
                level = "".join(ch for ch in style if ch.isdigit()) or "1"
                parts.append(f"{'#' * min(int(level), 6)} {text}")
            else:
                parts.append(text)
        for table in document.tables:
            rows = [[cell.text for cell in row.cells] for row in table.rows]
            parts.append(_rows_to_markdown(rows))
        return "\n\n".join(parts) if parts else "*(no readable text)*"


class XlsxExtractor:
    """Excel ``.xlsx`` → one Markdown table per sheet."""

    name = "xlsx"
    suffixes = (".xlsx",)

    def is_available(self) -> bool:
        return _has("openpyxl")

    def extract(self, path: Path) -> str:
        import openpyxl

        workbook = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        sections: list[str] = []
        for sheet in workbook.worksheets:
            rows = [list(row) for row in sheet.iter_rows(values_only=True)]
            sections.append(f"## {sheet.title}\n\n{_rows_to_markdown(rows)}")
        workbook.close()
        return "\n\n".join(sections) if sections else "*(no sheets)*"


class PptxExtractor:
    """PowerPoint ``.pptx`` → Markdown, one section per slide."""

    name = "pptx"
    suffixes = (".pptx",)

    def is_available(self) -> bool:
        return _has("pptx")

    def extract(self, path: Path) -> str:
        from pptx import Presentation

        presentation = Presentation(str(path))
        sections: list[str] = []
        for index, slide in enumerate(presentation.slides, start=1):
            lines = [f"## Slide {index}"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        lines.append(text)
            sections.append("\n\n".join(lines))
        return "\n\n".join(sections) if sections else "*(no slides)*"


class HtmlExtractor:
    """HTML → Markdown via ``markdownify`` (rich), or a stdlib text fallback."""

    name = "html"
    suffixes = (".html", ".htm")

    def is_available(self) -> bool:
        return True  # always usable — falls back to a stdlib text strip

    def extract(self, path: Path) -> str:
        html = path.read_text(encoding="utf-8", errors="replace")
        if _has("markdownify"):
            from markdownify import markdownify

            return markdownify(html, heading_style="ATX").strip()
        # Fallback: strip tags to readable text without a third-party dep.
        return _StripTags.strip(html)


class PdfExtractor:
    """PDF → text, page by page (``pypdf``). Mirrors the ``pdf`` tool's engine."""

    name = "pdf"
    suffixes = (".pdf",)

    def is_available(self) -> bool:
        return _has("pypdf")

    def extract(self, path: Path) -> str:
        from pypdf import PdfReader

        reader = PdfReader(str(path))
        pages = []
        for index, page in enumerate(reader.pages, start=1):
            text = (page.extract_text() or "").strip()
            pages.append(f"## Page {index}\n\n{text}" if text else f"## Page {index}\n\n*(no text)*")
        return "\n\n".join(pages) if pages else "*(no pages)*"


class _StripTags:
    """Minimal, dependency-free HTML→text fallback (drops script/style/tags)."""

    @staticmethod
    def strip(html: str) -> str:
        from html.parser import HTMLParser

        class _Parser(HTMLParser):
            def __init__(self) -> None:
                super().__init__()
                self._chunks: list[str] = []
                self._skip = 0

            def handle_starttag(self, tag: str, attrs: Any) -> None:
                if tag in ("script", "style"):
                    self._skip += 1
                elif tag in ("p", "br", "div", "li", "tr", "h1", "h2", "h3", "h4"):
                    self._chunks.append("\n")

            def handle_endtag(self, tag: str) -> None:
                if tag in ("script", "style") and self._skip:
                    self._skip -= 1

            def handle_data(self, data: str) -> None:
                if not self._skip and data.strip():
                    self._chunks.append(data)

        parser = _Parser()
        parser.feed(html)
        text = "".join(parser._chunks)
        return "\n".join(line.strip() for line in text.splitlines() if line.strip())


# Register the built-in extractors at import.
for _extractor in (
    CsvExtractor(),
    DocxExtractor(),
    XlsxExtractor(),
    PptxExtractor(),
    HtmlExtractor(),
    PdfExtractor(),
):
    register_extractor(_extractor)
