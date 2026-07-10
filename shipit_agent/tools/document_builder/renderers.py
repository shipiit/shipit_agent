"""Format renderers for DocumentBuilderTool.

Each renderer takes (title, sections|sheets, path) and writes the file.
All third-party imports are local so the base install stays dependency-free;
a missing library raises ImportError with the exact `pip install` fix.
"""

from __future__ import annotations

import html as _html
from pathlib import Path
from typing import Any

# Shared palette — one accent used consistently across formats.
ACCENT = "1F4E79"  # deep blue
ACCENT_LIGHT = "DCE6F1"


def _need(module: str, pip_name: str):
    try:
        return __import__(module)
    except ImportError as err:
        raise ImportError(
            f"This output needs the optional `{pip_name}` package — "
            f"install with `pip install {pip_name}`."
        ) from err


def render_pdf(title: str, sections: list[dict[str, Any]], path: Path) -> None:
    _need("reportlab", "reportlab")
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib.units import cm
    from reportlab.platypus import (
        Paragraph,
        SimpleDocTemplate,
        Spacer,
        Table,
        TableStyle,
    )

    styles = getSampleStyleSheet()
    styles["Title"].textColor = colors.HexColor(f"#{ACCENT}")
    styles["Heading2"].textColor = colors.HexColor(f"#{ACCENT}")
    story: list[Any] = [Paragraph(_html.escape(title), styles["Title"]), Spacer(1, 12)]

    for section in sections:
        if section.get("heading"):
            story.append(Paragraph(_html.escape(str(section["heading"])), styles["Heading2"]))
        if section.get("body"):
            for para in str(section["body"]).split("\n\n"):
                story.append(Paragraph(_html.escape(para), styles["BodyText"]))
        for bullet in section.get("bullets") or []:
            story.append(
                Paragraph(f"•  {_html.escape(str(bullet))}", styles["BodyText"])
            )
        table = section.get("table")
        if table and table.get("rows") is not None:
            data = [[str(c) for c in table.get("headers", [])]] if table.get("headers") else []
            data += [[str(c) for c in row] for row in table["rows"]]
            t = Table(data, hAlign="LEFT")
            t.setStyle(
                TableStyle(
                    [
                        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor(f"#{ACCENT}")),
                        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                        ("ROWBACKGROUNDS", (0, 1), (-1, -1),
                         [colors.white, colors.HexColor(f"#{ACCENT_LIGHT}")]),
                        ("GRID", (0, 0), (-1, -1), 0.4, colors.HexColor("#B0BEC5")),
                        ("FONTSIZE", (0, 0), (-1, -1), 9),
                        ("TOPPADDING", (0, 0), (-1, -1), 4),
                        ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
                    ]
                )
            )
            story.append(Spacer(1, 6))
            story.append(t)
        story.append(Spacer(1, 12))

    SimpleDocTemplate(
        str(path), pagesize=A4,
        leftMargin=2 * cm, rightMargin=2 * cm, topMargin=2 * cm, bottomMargin=2 * cm,
        title=title,
    ).build(story)


def render_xlsx(title: str, sheets: list[dict[str, Any]], path: Path) -> None:
    _need("openpyxl", "openpyxl")
    from openpyxl import Workbook
    from openpyxl.styles import Alignment, Font, PatternFill
    from openpyxl.utils import get_column_letter

    wb = Workbook()
    wb.remove(wb.active)
    header_font = Font(bold=True, color="FFFFFF")
    header_fill = PatternFill("solid", fgColor=ACCENT)

    for sheet in sheets or [{"name": "Sheet1", "headers": [], "rows": []}]:
        ws = wb.create_sheet(str(sheet.get("name", "Sheet"))[:31])
        headers = [str(h) for h in sheet.get("headers", [])]
        if headers:
            ws.append(headers)
            for cell in ws[1]:
                cell.font = header_font
                cell.fill = header_fill
                cell.alignment = Alignment(horizontal="center")
            ws.freeze_panes = "A2"
        for row in sheet.get("rows", []):
            ws.append(list(row))  # "=SUM(...)" strings become live formulas
        # Auto-size columns from content (formula text length is fine).
        for idx in range(1, ws.max_column + 1):
            width = max(
                (len(str(c.value)) for c in ws[get_column_letter(idx)] if c.value is not None),
                default=8,
            )
            ws.column_dimensions[get_column_letter(idx)].width = min(width + 3, 50)
    wb.properties.title = title
    wb.save(str(path))


def render_docx(title: str, sections: list[dict[str, Any]], path: Path) -> None:
    _need("docx", "python-docx")
    import docx

    doc = docx.Document()
    doc.add_heading(title, level=0)
    for section in sections:
        if section.get("heading"):
            doc.add_heading(str(section["heading"]), level=1)
        if section.get("body"):
            for para in str(section["body"]).split("\n\n"):
                doc.add_paragraph(para)
        for bullet in section.get("bullets") or []:
            doc.add_paragraph(str(bullet), style="List Bullet")
        table = section.get("table")
        if table and table.get("rows") is not None:
            headers = [str(h) for h in table.get("headers", [])]
            rows = table["rows"]
            cols = len(headers) or (len(rows[0]) if rows else 1)
            t = doc.add_table(rows=0, cols=cols)
            t.style = "Light Grid Accent 1"
            if headers:
                cells = t.add_row().cells
                for i, header in enumerate(headers):
                    run = cells[i].paragraphs[0].add_run(header)
                    run.bold = True
            for row in rows:
                cells = t.add_row().cells
                for i, value in enumerate(row):
                    cells[i].text = str(value)
    doc.save(str(path))


def render_pptx(title: str, sections: list[dict[str, Any]], path: Path) -> None:
    _need("pptx", "python-pptx")
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Pt

    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = title
    title_slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(ACCENT)

    for section in sections:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = str(section.get("heading", ""))
        slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(ACCENT)
        body = slide.placeholders[1].text_frame
        lines = list(section.get("bullets") or [])
        if section.get("body"):
            lines = [section["body"], *lines]
        for i, line in enumerate(lines):
            para = body.paragraphs[0] if i == 0 else body.add_paragraph()
            para.text = str(line)
            para.font.size = Pt(18)
    prs.save(str(path))


def render_html(title: str, sections: list[dict[str, Any]], path: Path) -> None:
    """Zero-dependency styled HTML — also a print-to-PDF fallback."""
    esc = _html.escape
    parts = [
        "<!doctype html><html><head><meta charset='utf-8'>",
        f"<title>{esc(title)}</title><style>",
        "body{font:15px/1.6 -apple-system,'Segoe UI',sans-serif;max-width:820px;"
        "margin:3rem auto;padding:0 1.5rem;color:#1a202c}",
        f"h1{{color:#{ACCENT};font-size:2rem}} h2{{color:#{ACCENT};margin-top:2rem}}",
        "table{border-collapse:collapse;width:100%;margin:1rem 0;font-size:.92em}",
        f"th{{background:#{ACCENT};color:#fff;text-align:left}}",
        "th,td{padding:.5rem .75rem;border:1px solid #cbd5e0}",
        f"tr:nth-child(even) td{{background:#{ACCENT_LIGHT}}}",
        "</style></head><body>",
        f"<h1>{esc(title)}</h1>",
    ]
    for section in sections:
        if section.get("heading"):
            parts.append(f"<h2>{esc(str(section['heading']))}</h2>")
        if section.get("body"):
            for para in str(section["body"]).split("\n\n"):
                parts.append(f"<p>{esc(para)}</p>")
        bullets = section.get("bullets") or []
        if bullets:
            items = "".join(f"<li>{esc(str(b))}</li>" for b in bullets)
            parts.append(f"<ul>{items}</ul>")
        table = section.get("table")
        if table and table.get("rows") is not None:
            head = "".join(f"<th>{esc(str(h))}</th>" for h in table.get("headers", []))
            body_rows = "".join(
                "<tr>" + "".join(f"<td>{esc(str(c))}</td>" for c in row) + "</tr>"
                for row in table["rows"]
            )
            parts.append(f"<table><tr>{head}</tr>{body_rows}</table>")
    parts.append("</body></html>")
    path.write_text("".join(parts), encoding="utf-8")
